#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
生成项目概览PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor as RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色主题
PRIMARY = RgbColor(0x1A, 0x5F, 0xB4)      # 深蓝
SECONDARY = RgbColor(0x2E, 0xCC, 0x71)    # 绿色
ACCENT = RgbColor(0xE7, 0x4C, 0x3C)       # 红色
DARK = RgbColor(0x2C, 0x3E, 0x50)         # 深色文字
LIGHT_BG = RgbColor(0xF8, 0xF9, 0xFA)     # 浅灰背景
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
ORANGE = RgbColor(0xE6, 0x7E, 0x22)
PURPLE = RgbColor(0x9B, 0x59, 0xB6)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    # 背景渐变效果 - 使用纯色矩形模拟
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = PRIMARY
    bg.line.fill.background()

    # 装饰圆形
    circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(-1), Inches(4), Inches(4))
    circle1.fill.solid()
    circle1.fill.fore_color.rgb = RgbColor(0x34, 0x7A, 0xD6)
    circle1.line.fill.background()
    circle1.rotation = 45

    circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-1.5), Inches(5), Inches(3), Inches(3))
    circle2.fill.solid()
    circle2.fill.fore_color.rgb = RgbColor(0x34, 0x7A, 0xD6)
    circle2.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 副标题
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0xBB, 0xD6, 0xF0)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(prs, title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # 装饰线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(4.3), Inches(3.333), Inches(0.05))
    line.fill.solid()
    line.fill.fore_color.rgb = SECONDARY
    line.line.fill.background()

    return slide

def add_content_slide(prs, title, bullets, color=PRIMARY):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 白色背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    # 顶部色条
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color

    # 内容
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(12), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {bullet}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK
        p.space_before = Pt(14)
        p.space_after = Pt(6)
        p.level = 0

    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, color=PRIMARY):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = color

    # 左列标题
    left_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    # 左列内容
    left_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(5.8), Inches(5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.6), Inches(1.3), Inches(0.01), Inches(5.2))
    line.fill.solid()
    line.fill.fore_color.rgb = RgbColor(0xDD, 0xDD, 0xDD)
    line.line.fill.background()

    # 右列标题
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(5.8), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    # 右列内容
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.8), Inches(5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"●  {item}"
        p.font.size = Pt(15)
        p.font.color.rgb = DARK
        p.space_before = Pt(8)
        p.space_after = Pt(4)

    return slide

def add_architecture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PURPLE
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "三层 POM 架构设计"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    # 架构图 - 三层堆叠
    layer_colors = [RgbColor(0x34, 0x98, 0xDB), RgbColor(0x2E, 0xCC, 0x71), RgbColor(0xE6, 0x7E, 0x22)]
    layer_titles = ["测试层 Test Layer", "POM 层 Page Object Layer", "基础层 Base Layer"]
    layer_subs = [
        "testcase/ 目录 - pytest测试用例",
        "pom/ 目录 - 页面对象封装",
        "bases/ 目录 - UI操作封装"
    ]
    layer_details = [
        ["16个测试文件", "覆盖7大功能模块", "Allure报告集成", "Fixture环境管理"],
        ["12个Page类", "控件链式定位", "图像识别定位", "业务方法封装"],
        ["BasePage基类 (949行)", "ScreenElement图像识别", "SlowMouseController", "Windows API调用"]
    ]

    y_start = 1.4
    layer_height = 1.6
    for i, (color, title, sub, details) in enumerate(zip(layer_colors, layer_titles, layer_subs, layer_details)):
        y = Inches(y_start + i * (layer_height + 0.2))

        # 层背景
        rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.5), y, Inches(10), Inches(layer_height))
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()

        # 层标题
        tbox = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.1), Inches(4), Inches(0.4))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE

        # 层副标题
        sbox = slide.shapes.add_textbox(Inches(1.8), y + Inches(0.45), Inches(4), Inches(0.3))
        tf = sbox.text_frame
        p = tf.paragraphs[0]
        p.text = sub
        p.font.size = Pt(13)
        p.font.color.rgb = RgbColor(0xEE, 0xEE, 0xEE)

        # 详情列表
        dbox = slide.shapes.add_textbox(Inches(6), y + Inches(0.15), Inches(5), Inches(1.3))
        tf = dbox.text_frame
        tf.word_wrap = True
        for j, d in enumerate(details):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {d}"
            p.font.size = Pt(13)
            p.font.color.rgb = WHITE
            p.space_before = Pt(4)

    # 箭头
    for i in range(2):
        y = Inches(y_start + (i + 1) * (layer_height + 0.2) - 0.15)
        arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(6.2), y, Inches(0.6), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RgbColor(0x95, 0x95, 0x95)
        arrow.line.fill.background()

    return slide

def add_code_stats_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RgbColor(0x16, 0xA0, 0x85)
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "代码规模统计"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x16, 0xA0, 0x85)

    # 统计卡片数据
    stats = [
        ("总代码行数", "~10,000+", "Python代码", RgbColor(0x34, 0x98, 0xDB)),
        ("POM页面", "12个", "Page类封装", RgbColor(0x2E, 0xCC, 0x71)),
        ("测试文件", "16个", "pytest测试", RgbColor(0xE6, 0x7E, 0x22)),
        ("工具模块", "17个", "commons/utils", RgbColor(0x9B, 0x59, 0xB6)),
    ]

    for i, (label, value, desc, color) in enumerate(stats):
        x = Inches(0.8 + i * 3.1)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.5), Inches(2.8), Inches(1.8))
        card.fill.solid()
        card.fill.fore_color.rgb = color
        card.line.fill.background()

        vbox = slide.shapes.add_textbox(x, Inches(1.7), Inches(2.8), Inches(0.7))
        tf = vbox.text_frame
        p = tf.paragraphs[0]
        p.text = value
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        lbox = slide.shapes.add_textbox(x, Inches(2.4), Inches(2.8), Inches(0.4))
        tf = lbox.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        dbox = slide.shapes.add_textbox(x, Inches(2.8), Inches(2.8), Inches(0.3))
        tf = dbox.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = RgbColor(0xEE, 0xEE, 0xEE)
        p.alignment = PP_ALIGN.CENTER

    # 各模块代码量详情
    details = [
        ("bases/basePage.py", "949行", "核心基类，封装所有UI操作"),
        ("commons/utils/configmanager.py", "1080行", "配置管理，INI文件操作"),
        ("bases/captureScreen.py", "363行", "图像识别与元素定位"),
        ("pom/voicewave_home_page.py", "590行", "首页/实时变声页面"),
        ("pom/voicewave_file_voicechanger_page.py", "474行", "文件变声页面"),
        ("testcase/test_nav4_file_voice_changer.py", "543行", "文件变声测试"),
        ("testcase/test_nav5_voice_creation.py", "424行", "声音创建测试"),
        ("testcase/test_language.py", "416行", "多语言测试"),
    ]

    detail_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.7), Inches(12), Inches(3.5))
    tf = detail_box.text_frame
    tf.word_wrap = True

    for i, (file, lines, desc) in enumerate(details):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{file:<45} {lines:>8}  →  {desc}"
        p.font.size = Pt(14)
        p.font.color.rgb = DARK
        p.font.name = "Consolas"
        p.space_before = Pt(6)
        p.space_after = Pt(3)

    return slide

def add_tech_stack_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RgbColor(0xE7, 0x4C, 0x3C)
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "核心技术栈"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0xE7, 0x4C, 0x3C)

    # 技术栈分类
    categories = [
        ("UI 自动化", [
            "uiautomation 2.0.29 - Windows UI自动化核心库",
            "pyautogui 0.9.54 - 鼠标键盘模拟",
            "pyperclip 1.8.2 - 剪贴板操作",
        ], RgbColor(0x34, 0x98, 0xDB)),
        ("图像识别", [
            "opencv-python 4.12.0.88 - OpenCV图像处理",
            "numpy 2.2.6 - 数值计算",
            "Pillow 11.1.0 - 图像处理",
        ], RgbColor(0x2E, 0xCC, 0x71)),
        ("测试框架", [
            "pytest 9.0.2 - 测试框架核心",
            "pytest-dependency 0.6.1 - 用例依赖管理",
            "pytest-html 4.2.0 - HTML报告",
            "allure-pytest 2.14.0 - Allure测试报告",
        ], RgbColor(0xE6, 0x7E, 0x22)),
        ("系统工具", [
            "WMI 1.5.1 - Windows管理接口",
            "psutil 7.2.1 - 进程管理",
            "selenium 4.41.0 - 浏览器自动化",
            "paramiko 3.5.1 - SSH远程操作",
        ], RgbColor(0x9B, 0x59, 0xB6)),
    ]

    for i, (cat_title, items, color) in enumerate(categories):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.4 + row * 2.8)

        # 分类卡片
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(2.5))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = color
        card.line.width = Pt(2)

        # 分类标题
        tbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.4))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = cat_title
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = color

        # 列表项
        ibox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.55), Inches(5.6), Inches(1.8))
        tf = ibox.text_frame
        tf.word_wrap = True
        for j, item in enumerate(items):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"▸ {item}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK
            p.space_before = Pt(6)

    return slide

def add_features_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = SECONDARY
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "框架核心特性"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = SECONDARY

    features = [
        ("双模式元素定位", "控件链式定位 + OpenCV图像识别，适应不同场景", RgbColor(0x34, 0x98, 0xDB)),
        ("拟人化操作", "平滑鼠标移动、缓慢点击，模拟真实用户行为", RgbColor(0xE7, 0x4C, 0x3C)),
        ("多语言支持", "支持12种语言测试：英/法/德/意/韩/葡/日/繁中/简中/西/土/阿", RgbColor(0x2E, 0xCC, 0x71)),
        ("用户环境隔离", "新用户/老用户/未激活用户环境，通过Fixture独立管理", RgbColor(0xE6, 0x7E, 0x22)),
        ("Allure报告", "自动生成可视化测试报告，支持截图附件", RgbColor(0x9B, 0x59, 0xB6)),
        ("滚动查找", "支持上下滚动查找控件和图像元素", RgbColor(0x16, 0xA0, 0x85)),
        ("文件管理器", "封装文件选择、路径输入、文件打开等操作", RgbColor(0x34, 0x98, 0xDB)),
        ("异常处理", "element_raise / control_raise 统一异常策略", RgbColor(0xE7, 0x4C, 0x3C)),
    ]

    for i, (feat, desc, color) in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(0.6 + col * 6.3)
        y = Inches(1.4 + row * 1.35)

        # 色块标记
        marker = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.05), Inches(0.08), Inches(0.9))
        marker.fill.solid()
        marker.fill.fore_color.rgb = color
        marker.line.fill.background()

        # 特性名
        fbox = slide.shapes.add_textbox(x + Inches(0.2), y, Inches(5.8), Inches(0.35))
        tf = fbox.text_frame
        p = tf.paragraphs[0]
        p.text = feat
        p.font.size = Pt(17)
        p.font.bold = True
        p.font.color.rgb = color

        # 描述
        dbox = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.35), Inches(5.8), Inches(0.6))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(13)
        p.font.color.rgb = DARK

    return slide

def add_pages_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = ORANGE
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "POM 页面覆盖"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ORANGE

    pages = [
        ("voicewave_home_page", "590行", "首页 / 实时变声", "核心功能页面"),
        ("voicewave_file_voicechanger_page", "474行", "文件变声", "音频文件处理"),
        ("voicewave_setting_page", "269行", "设置页面", "程序配置"),
        ("voicewave_discount_alert", "250行", "弹窗处理", "启动弹窗/折扣弹窗"),
        ("voicewave_soundboard_page", "167行", "音效板", "Soundboard功能"),
        ("voicewave_voice_creation_page", "168行", "声音创建", "Voice Creation"),
        ("voicewave_community_page", "173行", "社区", "Community Library"),
        ("voicewave_language_page", "155行", "语言设置", "多语言切换"),
        ("voicewave_closeprogram_page", "126行", "关闭程序", "退出逻辑"),
        ("voicewave_login_page", "123行", "登录页面", "用户登录"),
        ("voicewave_active", "38行", "激活管理", "激活状态"),
        ("voicewave_voice_page", "20行", "声音页面", "Voice基础"),
    ]

    for i, (name, lines, func, note) in enumerate(pages):
        col = i % 3
        row = i // 3
        x = Inches(0.5 + col * 4.2)
        y = Inches(1.4 + row * 1.4)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(4), Inches(1.25))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RgbColor(0xDD, 0xDD, 0xDD)

        nbox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.1), Inches(3.7), Inches(0.35))
        tf = nbox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.font.name = "Consolas"

        fbox = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.42), Inches(2.5), Inches(0.3))
        tf = fbox.text_frame
        p = tf.paragraphs[0]
        p.text = func
        p.font.size = Pt(12)
        p.font.color.rgb = ORANGE

        lbox = slide.shapes.add_textbox(x + Inches(2.8), y + Inches(0.42), Inches(1.1), Inches(0.3))
        tf = lbox.text_frame
        p = tf.paragraphs[0]
        p.text = lines
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(0x95, 0x95, 0x95)
        p.alignment = PP_ALIGN.RIGHT

        note_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.75), Inches(3.7), Inches(0.3))
        tf = note_box.text_frame
        p = tf.paragraphs[0]
        p.text = note
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(0x95, 0x95, 0x95)

    return slide

def add_testcases_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RgbColor(0x16, 0xA0, 0x85)
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "测试用例覆盖"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x16, 0xA0, 0x85)

    testcases = [
        ("test_nav4_file_voice_changer", "543行", "文件变声功能测试", "核心"),
        ("test_nav5_voice_creation", "424行", "声音创建功能测试", "核心"),
        ("test_language", "416行", "多语言UI测试", "重要"),
        ("test_nav3_community", "285行", "社区功能测试", "核心"),
        ("test_title", "256行", "窗口标题测试", "基础"),
        ("test_nav2_soundboard", "243行", "音效板功能测试", "核心"),
        ("test_about", "145行", "关于页面测试", "基础"),
        ("test_login", "177行", "登录功能测试", "重要"),
        ("test_testing", "127行", "通用测试方法", "基础"),
        ("test_left_navigation_bar", "115行", "左侧导航栏测试", "基础"),
        ("test_inactive", "115行", "未激活用户测试", "重要"),
        ("test_olduser_start_alert", "63行", "老用户启动弹窗", "基础"),
        ("test_newuser_start_alert", "56行", "新用户启动弹窗", "基础"),
    ]

    # 表格方式展示
    table_left = Inches(0.6)
    table_top = Inches(1.4)
    table_width = Inches(12)
    table_height = Inches(5.5)

    rows, cols = len(testcases) + 1, 4
    table = slide.shapes.add_table(rows, cols, table_left, table_top, table_width, table_height).table

    # 表头
    headers = ["测试文件", "代码行", "测试内容", "优先级"]
    header_colors = [RgbColor(0x16, 0xA0, 0x85)] * 4
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_colors[i]

    # 数据行
    for row_idx, (name, lines, content, priority) in enumerate(testcases):
        row = row_idx + 1
        data = [name, lines, content, priority]
        colors = [DARK, RgbColor(0x95, 0x95, 0x95), DARK,
                   RgbColor(0xE7, 0x4C, 0x3C) if priority == "核心" else (
                   RgbColor(0xE6, 0x7E, 0x22) if priority == "重要" else RgbColor(0x16, 0xA0, 0x85))]
        for col_idx, (val, color) in enumerate(zip(data, colors)):
            cell = table.cell(row, col_idx)
            cell.text = val
            cell.text_frame.paragraphs[0].font.size = Pt(12)
            cell.text_frame.paragraphs[0].font.color.rgb = color
            if col_idx == 0:
                cell.text_frame.paragraphs[0].font.name = "Consolas"
            # 交替行背景
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return slide

def add_fixture_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = PURPLE
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Fixture 环境管理"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = PURPLE

    fixtures = [
        ("active_window", "module", "已激活用户环境", "SHOW_EXIT_WINDOW=false"),
        ("newuser_language_window", "function", "未激活新用户 + 语言切换", "delete_voice_wave_ini()"),
        ("olduser_language_window", "function", "未激活老用户 + 语言切换", "modify_time() + delete_start_time()"),
        ("creation_language_window", "function", "声音创建多语言环境", "6种语言参数化"),
        ("login_window", "function", "登录测试环境", "delete_voice_wave_ini()"),
        ("active_window_function", "function", "已激活用户(含清理)", "完整前后置"),
        ("window", "module", "标准窗口(含清理)", "delete_voice_wave_ini()"),
        ("bottom_window", "function", "底部退出测试环境", "测试退出流程"),
        ("main_window", "class", "新用户完整环境", "含弹窗关闭 + 进程清理"),
        ("old_main_window", "function", "老用户完整环境", "modify_time()"),
        ("new_main_window", "function", "新用户环境", "delete_voice_wave_ini()"),
        ("Inactive_main_window", "function", "未激活环境", "清除激活文件"),
    ]

    for i, (name, scope, desc, setup) in enumerate(fixtures):
        col = i % 2
        row = i // 2
        x = Inches(0.5 + col * 6.3)
        y = Inches(1.35 + row * 0.92)

        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(6), Inches(0.82))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_BG
        card.line.color.rgb = RgbColor(0xDD, 0xDD, 0xDD)

        # 作用域标签
        scope_colors = {
            "session": RgbColor(0x9B, 0x59, 0xB6),
            "module": RgbColor(0x34, 0x98, 0xDB),
            "class": RgbColor(0xE6, 0x7E, 0x22),
            "function": RgbColor(0x2E, 0xCC, 0x71),
        }
        tag = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), y + Inches(0.1), Inches(0.9), Inches(0.25))
        tag.fill.solid()
        tag.fill.fore_color.rgb = scope_colors.get(scope, RgbColor(0x95, 0x95, 0x95))
        tag.line.fill.background()

        tbox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.08), Inches(0.9), Inches(0.25))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = scope
        p.font.size = Pt(10)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 名称
        nbox = slide.shapes.add_textbox(x + Inches(1.1), y + Inches(0.08), Inches(4.7), Inches(0.28))
        tf = nbox.text_frame
        p = tf.paragraphs[0]
        p.text = name
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.font.name = "Consolas"

        # 描述
        dbox = slide.shapes.add_textbox(x + Inches(0.1), y + Inches(0.4), Inches(5.8), Inches(0.3))
        tf = dbox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{desc}  |  Setup: {setup}"
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(0x77, 0x77, 0x77)

    return slide

def add_workflow_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = RgbColor(0x34, 0x98, 0xDB)
    top_bar.line.fill.background()

    # 标题
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "测试执行流程"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RgbColor(0x34, 0x98, 0xDB)

    # 流程步骤
    steps = [
        ("1", "环境准备", "ConfigManager设置INI配置\n删除/修改配置文件\n注册表切换语言"),
        ("2", "启动应用", "subprocess.Popen启动程序\n等待主窗口出现(最多120秒)\n获取WindowControl对象"),
        ("3", "执行测试", "POM页面对象调用业务方法\n控件定位 / 图像识别\n断言验证结果"),
        ("4", "截图报告", "Allure截图附加到报告\n日志记录到文件\n生成allure-results"),
        ("5", "环境清理", "关闭程序弹窗处理\n选择退出方式\nkill_process清理进程"),
    ]

    for i, (num, step_title, step_desc) in enumerate(steps):
        x = Inches(0.8 + i * 2.4)
        y = Inches(1.8)

        # 步骤编号圆
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.7), Inches(1.3), Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = RgbColor(0x34, 0x98, 0xDB)
        circle.line.fill.background()

        nbox = slide.shapes.add_textbox(x + Inches(0.7), Inches(1.35), Inches(0.6), Inches(0.5))
        tf = nbox.text_frame
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # 标题
        tbox = slide.shapes.add_textbox(x, Inches(2.1), Inches(2.4), Inches(0.4))
        tf = tbox.text_frame
        p = tf.paragraphs[0]
        p.text = step_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK
        p.alignment = PP_ALIGN.CENTER

        # 描述
        dbox = slide.shapes.add_textbox(x, Inches(2.5), Inches(2.4), Inches(2))
        tf = dbox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step_desc
        p.font.size = Pt(11)
        p.font.color.rgb = RgbColor(0x77, 0x77, 0x77)
        p.alignment = PP_ALIGN.CENTER

        # 箭头
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.2), Inches(1.55), Inches(0.35), Inches(0.2))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RgbColor(0xBB, 0xBB, 0xBB)
            arrow.line.fill.background()

    # 底部命令
    cmd_box = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(12), Inches(1.8))
    tf = cmd_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "常用命令："
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = DARK

    commands = [
        "pytest -v -s -m test              # 运行所有标记为test的用例",
        "pytest -m \"ui\" -v               # 运行UI自动化测试",
        "pytest --alluredir=allure-results --clean-alluredir   # 生成Allure数据",
        "allure generate -o report -c allure-results           # 生成可视化报告",
    ]
    for cmd in commands:
        p = tf.add_paragraph()
        p.text = f"  $ {cmd}"
        p.font.size = Pt(13)
        p.font.color.rgb = RgbColor(0x16, 0xA0, 0x85)
        p.font.name = "Consolas"
        p.space_before = Pt(6)

    return slide

def add_end_slide(prs):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    # 装饰
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.2), Inches(1.8), Inches(3), Inches(3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RgbColor(0x34, 0x49, 0x5E)
    circle.line.fill.background()

    # 感谢文字
    tbox = slide.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.3), Inches(1))
    tf = tbox.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢聆听"
    p.font.size = Pt(56)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    sbox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.3), Inches(0.6))
    tf = sbox.text_frame
    p = tf.paragraphs[0]
    p.text = "EaseUS VoiceWave UI 自动化测试框架"
    p.font.size = Pt(22)
    p.font.color.rgb = RgbColor(0xBB, 0xBB, 0xBB)
    p.alignment = PP_ALIGN.CENTER

    return slide


# ==================== 生成所有幻灯片 ====================

# 1. 封面
add_title_slide(prs,
    "EaseUS VoiceWave",
    "PC端 UI 自动化测试框架项目概览")

# 2. 章节：项目概述
add_section_slide(prs, "项目概述")

# 3. 项目简介
add_content_slide(prs, "项目简介", [
    "目标应用：EaseUS VoiceWave — Windows桌面端音频变声软件",
    "技术方案：Python + uiautomation + pytest + Allure",
    "架构模式：三层 POM（Page Object Model）页面对象模型",
    "总代码量：约 10,000+ 行 Python 代码",
    "覆盖模块：7大功能模块（实时变声、文件变声、音效板、社区、声音创建、设置、登录）",
    "测试语言：支持 12 种语言（英语、法语、德语、意大利语、韩语、葡萄牙语、日语、繁中、简中、西班牙语、土耳其语、阿拉伯语）",
], PRIMARY)

# 4. 核心技术栈
add_tech_stack_slide(prs)

# 5. 章节：架构设计
add_section_slide(prs, "架构设计")

# 6. 三层POM架构
add_architecture_slide(prs)

# 7. 代码规模统计
add_code_stats_slide(prs)

# 8. 框架核心特性
add_features_slide(prs)

# 9. 章节：页面与用例
add_section_slide(prs, "页面与用例")

# 10. POM页面覆盖
add_pages_slide(prs)

# 11. 测试用例覆盖
add_testcases_slide(prs)

# 12. Fixture环境管理
add_fixture_slide(prs)

# 13. 章节：执行与报告
add_section_slide(prs, "执行与报告")

# 14. 测试执行流程
add_workflow_slide(prs)

# 15. 结束页
add_end_slide(prs)

# 保存
output_path = r"G:\project\python-uiautomation-for-pc-evw\项目概览_EaseUS_VoiceWave_UI自动化测试框架.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
print(f"共 {len(prs.slides)} 页幻灯片")
